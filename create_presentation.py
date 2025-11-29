from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Helper to add a slide with title and content
    def add_slide(title_text, content_text_list):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = title_text
        
        tf = content.text_frame
        tf.clear()
        
        for i, text in enumerate(content_text_list):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = text
            p.level = 0
            p.font.size = Pt(18)

    # Helper to add a section header slide
    def add_section_slide(title_text):
        slide_layout = prs.slide_layouts[2] # Section Header
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text

    # 1. Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "BootCampMusic Project"
    subtitle.text = "Comprehensive System Overview & Development Report\nWaterfall Methodology"

    # 2. Project Overview
    add_section_slide("1. Project Overview")
    add_slide("Project Objectives", [
        "Develop a full-stack music streaming platform.",
        "Provide seamless user experience with instant playback.",
        "Enable administrative management of music assets.",
        "Ensure scalability for future cloud deployment."
    ])
    add_slide("Project Scope", [
        "Frontend: Responsive Web Application (React).",
        "Backend: RESTful API (Django).",
        "Database: Relational Database for structured data.",
        "Features: Music Playback, Search, Admin Dashboard, Authentication."
    ])

    # 3. Requirement Analysis
    add_section_slide("2. Requirement Analysis")
    add_slide("Functional Requirements", [
        "User Authentication: Sign up, Login, Logout.",
        "Music Discovery: Search by title, artist, genre.",
        "Playback: Global player, Hover-to-preview.",
        "Admin: Upload, Edit, Delete tracks/artists/albums.",
        "Download: Secure file download for authenticated users."
    ])
    add_slide("Non-Functional Requirements", [
        "Performance: Instant playback start (<200ms).",
        "Usability: Intuitive UI with dark mode aesthetics.",
        "Reliability: Error handling for missing files or network issues.",
        "Maintainability: Modular code structure (Component-based)."
    ])

    # 4. System Design
    add_section_slide("3. System Design")
    add_slide("System Architecture", [
        "Client-Server Architecture.",
        "Frontend: React (SPA) communicating via Axios.",
        "Backend: Django REST Framework serving JSON APIs.",
        "Media Storage: Local file system (Dev) -> S3 (Prod).",
        "Database: SQLite (Dev) -> PostgreSQL (Prod)."
    ])
    add_slide("Technology Stack", [
        "Frontend: React, TypeScript, Tailwind CSS, Framer Motion.",
        "Backend: Python, Django, Django REST Framework.",
        "Database: SQLite (Development).",
        "Tools: VS Code, Git, Docker (Planned)."
    ])

    # 5. Database Design
    add_section_slide("4. Database Design")
    add_slide("Entity Relationship Diagram (ERD)", [
        "Artist: Name, Bio, Image.",
        "Album: Title, Release Date, Cover Image, FK(Artist).",
        "Track: Title, Duration, File, Preview, Genre, FK(Artist), FK(Album).",
        "DownloadLog: User, Track, Timestamp."
    ])

    # 6. Detailed Design
    add_section_slide("5. Detailed Design")
    add_slide("Frontend Components", [
        "Navbar: Navigation & Auth state management.",
        "TrackCard: Individual track display with Hover-to-Play logic.",
        "GlobalPlayer: Persistent footer player with progress/volume control.",
        "Pages: Home, Search, Login, Dashboard, Admin(Upload/Manage)."
    ])
    add_slide("Backend API Endpoints", [
        "GET /api/music/tracks/: List & Filter tracks.",
        "GET /api/music/tracks/{id}/: Track details.",
        "POST /api/admin/upload-track/: Admin upload (Multi-part).",
        "GET /api/auth/me/: Current user context."
    ])

    # 7. Implementation Highlights
    add_section_slide("6. Implementation Highlights")
    add_slide("Key Feature: Instant Playback", [
        "Challenge: Latency in starting audio.",
        "Solution: Preload='auto' on audio elements.",
        "Optimization: Reduced hover delay to 150ms.",
        "Result: Snappy, responsive user experience."
    ])
    add_slide("Key Feature: Admin Workflow", [
        "Complex dependency handling (Artist -> Album -> Track).",
        "Unified Upload Form: Creates Artist/Album automatically if missing.",
        "Transactional integrity ensuring data consistency."
    ])

    # 8. Testing & Status
    add_section_slide("7. Testing & Status")
    add_slide("Current Status", [
        "Development Phase: Completed.",
        "Local Testing: Passed (Functional & UI).",
        "Known Issues: Resolved (Playback delay, Admin bugs).",
        "Ready for: Containerization & Deployment."
    ])

    # 9. Future Plans (Deployment)
    add_section_slide("8. Future Plans")
    add_slide("Deployment Strategy", [
        "Containerization: Dockerize Frontend & Backend.",
        "Orchestration: Docker Compose for local simulation.",
        "Cloud Provider: AWS (EC2 or ECS).",
        "CI/CD: GitHub Actions for automated testing/deployment."
    ])

    # Save
    prs.save('BootCampMusic_Project_Presentation.pptx')
    print("Presentation saved successfully.")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("Error: python-pptx library is not installed.")
    except Exception as e:
        print(f"Error creating presentation: {e}")
